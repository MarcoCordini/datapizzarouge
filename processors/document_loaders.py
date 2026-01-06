"""
Document loaders per vari formati.
Estrae testo da PDF, Word, Excel, PowerPoint, Immagini (OCR).
"""
import logging
from pathlib import Path
from typing import Dict, List, Optional
import filetype

# PDF
import fitz  # PyMuPDF
import pdfplumber

# Office
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

# OCR
from PIL import Image
import easyocr

# Utilities
import re

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Carica e estrae testo da vari formati di documento."""

    def __init__(self, ocr_language: str = "it+en"):
        """
        Args:
            ocr_language: Lingue per OCR EasyOCR (es: "it", "en", "it+en")
                         Codici supportati: it, en, fr, de, es, etc.
        """
        self.ocr_language = ocr_language

        # Inizializza EasyOCR reader
        logger.info(f"Inizializzazione EasyOCR reader per lingue: {ocr_language}")
        languages = ocr_language.split("+")
        self.ocr_reader = easyocr.Reader(languages, gpu=False)
        logger.info("EasyOCR reader pronto")

    def load(self, file_path: str) -> Dict:
        """
        Carica documento e estrae testo.

        Returns:
            {
                "text": "Contenuto estratto...",
                "metadata": {
                    "file_name": "doc.pdf",
                    "file_type": "pdf",
                    "pages": 10,
                    "source": "/path/to/doc.pdf"
                }
            }
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File non trovato: {file_path}")

        # Auto-detect file type
        kind = filetype.guess(str(path))
        extension = path.suffix.lower()

        logger.info(f"Caricamento {path.name} (tipo: {extension})")

        # Router per tipo file
        if extension == ".pdf":
            return self._load_pdf(path)
        elif extension == ".docx":
            return self._load_docx(path)
        elif extension == ".doc":
            return self._load_doc_legacy(path)
        elif extension == ".pptx":
            return self._load_pptx(path)
        elif extension == ".xlsx":
            return self._load_xlsx(path)
        elif extension in [".txt", ".md", ".csv"]:
            return self._load_text(path)
        elif extension in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
            return self._load_image_ocr(path)
        else:
            raise ValueError(f"Formato non supportato: {extension}")

    # === PDF ===

    def _load_pdf(self, path: Path) -> Dict:
        """Carica PDF (testo o scansionato con OCR)."""
        try:
            # Prima prova estrazione testo nativo
            doc = fitz.open(str(path))
            text_parts = []
            has_text = False

            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                if text.strip():
                    has_text = True
                    text_parts.append(text)

            doc.close()

            # Se PDF ha testo nativo, usa quello
            if has_text:
                logger.info(f"PDF con testo nativo ({len(text_parts)} pagine)")
                return {
                    "text": "\n\n".join(text_parts),
                    "metadata": {
                        "file_name": path.name,
                        "file_type": "pdf",
                        "pages": len(text_parts),
                        "source": str(path.absolute()),
                        "extraction_method": "native_text"
                    }
                }

            # Altrimenti, usa OCR
            logger.info(f"PDF scansionato, uso OCR...")
            return self._load_pdf_ocr(path)

        except Exception as e:
            logger.error(f"Errore caricamento PDF: {e}")
            raise

    def _load_pdf_ocr(self, path: Path) -> Dict:
        """Carica PDF scansionato con OCR EasyOCR."""
        doc = fitz.open(str(path))
        text_parts = []

        for page_num, page in enumerate(doc, start=1):
            logger.info(f"OCR pagina {page_num}/{len(doc)}")

            # Renderizza pagina come immagine
            pix = page.get_pixmap(dpi=300)  # Alta risoluzione per OCR
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            # OCR con EasyOCR
            result = self.ocr_reader.readtext(img, detail=0)  # detail=0 per avere solo il testo
            text = "\n".join(result)
            text_parts.append(text)

        doc.close()

        return {
            "text": "\n\n".join(text_parts),
            "metadata": {
                "file_name": path.name,
                "file_type": "pdf",
                "pages": len(text_parts),
                "source": str(path.absolute()),
                "extraction_method": "easyocr"
            }
        }

    # === WORD ===

    def _load_docx(self, path: Path) -> Dict:
        """Carica Word DOCX con estrazione immagini."""
        doc = DocxDocument(str(path))

        # Estrai paragrafi
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]

        # Estrai tabelle
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        # Estrai immagini embedded
        images = self._extract_images_from_docx(doc, path)

        return {
            "text": "\n\n".join(text_parts),
            "metadata": {
                "file_name": path.name,
                "file_type": "docx",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "source": str(path.absolute()),
                "images": images  # Lista di informazioni sulle immagini
            }
        }

    def _extract_images_from_docx(self, doc: DocxDocument, doc_path: Path) -> list:
        """
        Estrae informazioni sulle immagini embedded in un file DOCX.

        Returns:
            Lista di dict con info immagini: [{
                "image_id": "rId5",
                "content_type": "image/png",
                "data": <bytes>,
                "paragraph_index": 3
            }]
        """
        images = []

        try:
            # Mappa delle relazioni (rId -> image data)
            image_rels = {}
            for rel_id, rel in doc.part.rels.items():
                if "image" in rel.target_ref:
                    image_rels[rel_id] = {
                        "content_type": rel.target_part.content_type,
                        "data": rel.target_part.blob
                    }

            # Trova posizione immagini nei paragrafi
            for para_idx, paragraph in enumerate(doc.paragraphs):
                for run in paragraph.runs:
                    # Cerca inline shapes (immagini)
                    if hasattr(run._element, 'drawing_lst'):
                        for drawing in run._element.drawing_lst:
                            # Cerca relationship id dell'immagine
                            blip = drawing.find('.//{*}blip')
                            if blip is not None:
                                embed_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if embed_id and embed_id in image_rels:
                                    images.append({
                                        "image_id": embed_id,
                                        "content_type": image_rels[embed_id]["content_type"],
                                        "data": image_rels[embed_id]["data"],
                                        "paragraph_index": para_idx,
                                        "text_before": paragraph.text[:100] if paragraph.text else ""
                                    })

            logger.info(f"Estratte {len(images)} immagini da {doc_path.name}")

        except Exception as e:
            logger.warning(f"Errore estrazione immagini da {doc_path.name}: {e}")

        return images

    def _load_doc_legacy(self, path: Path) -> Dict:
        """Carica Word DOC legacy (richiede LibreOffice o antiword)."""
        import subprocess

        # Prova con antiword (se installato)
        try:
            result = subprocess.run(
                ["antiword", str(path)],
                capture_output=True,
                text=True,
                check=True
            )
            text = result.stdout

            return {
                "text": text,
                "metadata": {
                    "file_name": path.name,
                    "file_type": "doc",
                    "source": str(path.absolute()),
                    "extraction_method": "antiword"
                }
            }
        except (FileNotFoundError, subprocess.CalledProcessError):
            raise ValueError(
                "Formato .doc legacy non supportato. "
                "Installa antiword o converti in .docx"
            )

    # === POWERPOINT ===

    def _load_pptx(self, path: Path) -> Dict:
        """Carica PowerPoint PPTX."""
        prs = Presentation(str(path))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []

            # Estrai testo da tutti gli shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)

            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        return {
            "text": "\n\n".join(text_parts),
            "metadata": {
                "file_name": path.name,
                "file_type": "pptx",
                "slides": len(prs.slides),
                "source": str(path.absolute())
            }
        }

    # === EXCEL ===

    def _load_xlsx(self, path: Path) -> Dict:
        """Carica Excel XLSX."""
        wb = openpyxl.load_workbook(str(path), data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_texts = [f"[Foglio: {sheet_name}]"]

            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell else "" for cell in row)
                if row_text.strip(" |"):
                    sheet_texts.append(row_text)

            if len(sheet_texts) > 1:  # Ha contenuto oltre al titolo
                text_parts.append("\n".join(sheet_texts))

        return {
            "text": "\n\n".join(text_parts),
            "metadata": {
                "file_name": path.name,
                "file_type": "xlsx",
                "sheets": len(wb.sheetnames),
                "source": str(path.absolute())
            }
        }

    # === TESTO ===

    def _load_text(self, path: Path) -> Dict:
        """Carica file di testo (TXT, MD, CSV)."""
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()

        return {
            "text": text,
            "metadata": {
                "file_name": path.name,
                "file_type": path.suffix.lower(),
                "source": str(path.absolute())
            }
        }

    # === IMMAGINI (OCR) ===

    def _load_image_ocr(self, path: Path) -> Dict:
        """Carica immagine con OCR EasyOCR."""
        logger.info(f"OCR su immagine {path.name}")

        img = Image.open(str(path))

        # OCR con EasyOCR
        result = self.ocr_reader.readtext(str(path), detail=0)  # detail=0 per avere solo il testo
        text = "\n".join(result)

        return {
            "text": text,
            "metadata": {
                "file_name": path.name,
                "file_type": "image",
                "dimensions": f"{img.width}x{img.height}",
                "source": str(path.absolute()),
                "extraction_method": "easyocr"
            }
        }


# === BATCH PROCESSING ===

class DocumentBatchLoader:
    """Carica batch di documenti da una directory."""

    def __init__(self, ocr_language: str = "it+en"):
        self.loader = DocumentLoader(ocr_language=ocr_language)

    def load_directory(
        self,
        directory: str,
        recursive: bool = True,
        extensions: Optional[List[str]] = None
    ) -> List[Dict]:
        """
        Carica tutti i documenti da una directory.

        Args:
            directory: Path directory
            recursive: Se True, cerca anche in subdirectory
            extensions: Lista estensioni da processare (es: [".pdf", ".docx"])
                        Se None, processa tutti i formati supportati

        Returns:
            Lista di documenti estratti
        """
        path = Path(directory)

        if not path.exists():
            raise FileNotFoundError(f"Directory non trovata: {directory}")

        if extensions is None:
            extensions = [
                ".pdf", ".docx", ".doc", ".pptx", ".xlsx",
                ".txt", ".md", ".jpg", ".jpeg", ".png"
            ]

        # Trova file
        if recursive:
            files = [f for f in path.rglob("*") if f.suffix.lower() in extensions]
        else:
            files = [f for f in path.glob("*") if f.suffix.lower() in extensions]

        logger.info(f"Trovati {len(files)} documenti in {directory}")

        # Carica
        documents = []
        for file_path in files:
            try:
                doc = self.loader.load(str(file_path))
                documents.append(doc)
                logger.info(f"OK - {file_path.name}")
            except Exception as e:
                logger.error(f"ERRORE caricando {file_path.name}: {e}")

        return documents
